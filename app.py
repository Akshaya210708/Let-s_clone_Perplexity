import os
import sqlite3
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from google import genai
from google.genai import types
from pypdf import PdfReader
from pptx import Presentation
from datetime import datetime
import io

app = Flask(__name__, static_folder='eco', static_url_path='')
CORS(app)

# Initialize Google GenAI Client with NEW API KEY
API_KEY = "AIzaSyAEi8BFaN2RW76x3ry-OhG2rNl7IC1QaAw"
PRIMARY_MODEL = "gemini-2.5-flash"
FALLBACK_MODEL = "gemini-1.5-flash"
DATABASE = 'history.db'

client = genai.Client(api_key=API_KEY)

# 1. Database Initialization
def init_db():
    with sqlite3.connect(DATABASE) as conn:
        # Create table if not exists with all columns
        conn.execute('''
            CREATE TABLE IF NOT EXISTS history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                query TEXT NOT NULL,
                response TEXT NOT NULL,
                filename TEXT,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        
        # Migration: Ensure 'filename' column exists even if table was created previously
        try:
            cursor = conn.cursor()
            cursor.execute('SELECT filename FROM history LIMIT 1')
        except sqlite3.OperationalError:
            print("Migration: Adding 'filename' column to history table.")
            conn.execute('ALTER TABLE history ADD COLUMN filename TEXT')
        
        conn.commit()

init_db()

def extract_text_from_pptx(file_bytes):
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        raise Exception(f"PowerPoint extraction failed: {str(e)}")

@app.route("/")
def index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route("/api/upload", methods=["POST"])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    try:
        file_bytes = file.read()
        filename = file.filename.lower()
        text = ""
        
        if filename.endswith('.pdf'):
            pdf_stream = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_stream)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        elif filename.endswith(('.pptx', '.ppt')):
            text = extract_text_from_pptx(file_bytes)
        else:
            return jsonify({"error": "Invalid file type. Only PDFs and PowerPoint are supported."}), 400
            
        return jsonify({
            "filename": file.filename,
            "text": text
        })
    except Exception as e:
        return jsonify({"error": f"Failed to parse file: {str(e)}"}), 500

@app.route("/api/chat", methods=["POST"])
def chat():
    data = request.json
    user_message = data.get("message")
    chat_history = data.get("history", []) 
    pdf_context = data.get("pdf_context", "") # Now only applies to the current message
    current_filename = data.get("filename", None)
    
    if not user_message:
        return jsonify({"error": "No message provided"}), 400
    
    # Construct System Instruction
    system_instruction = "You are a helpful and premium AI assistant named Eco. "
    if pdf_context:
        system_instruction += f"\n\nCONTEXT FROM UPLOADED FILE ({current_filename}):\n{pdf_context}\n\nPlease prioritize using this context to answer the specific query it is attached to."

    # Convert frontend history to Gemini format
    contents = []
    for msg in chat_history:
        contents.append(types.Content(
            role=msg['role'],
            parts=[types.Part(text=msg['text'])]
        ))
    
    # Append current message
    contents.append(types.Content(
        role='user',
        parts=[types.Part(text=user_message)]
    ))

    # Try primary model first, then fallback
    models_to_try = [PRIMARY_MODEL, FALLBACK_MODEL]
    
    for current_model in models_to_try:
        try:
            response = client.models.generate_content(
                model=current_model,
                contents=contents,
                config=types.GenerateContentConfig(
                    system_instruction=system_instruction
                )
            )
            
            reply_text = response.text
            
            # Save to History (Including the filename)
            with sqlite3.connect(DATABASE) as conn:
                cursor = conn.cursor()
                cursor.execute('''
                    INSERT INTO history (query, response, filename) 
                    VALUES (?, ?, ?)
                ''', (user_message, reply_text, current_filename))
                history_id = cursor.lastrowid
                conn.commit()
            
            return jsonify({
                "id": history_id,
                "reply": reply_text,
                "model": current_model
            })
            
        except Exception as e:
            error_msg = str(e)
            print(f"Error calling {current_model}: {error_msg}")
            
            if "RESOURCE_EXHAUSTED" in error_msg or "429" in error_msg:
                if current_model == PRIMARY_MODEL:
                    continue 
                else:
                    return jsonify({"error": "All models exhausted quota.", "status": "quota_exceeded"}), 429
            
            return jsonify({"error": error_msg}), 500

@app.route("/api/history", methods=["GET"])
def get_history():
    try:
        with sqlite3.connect(DATABASE) as conn:
            conn.row_factory = sqlite3.Row
            cursor = conn.cursor()
            cursor.execute('SELECT * FROM history ORDER BY timestamp DESC LIMIT 30')
            rows = cursor.fetchall()
            return jsonify([dict(row) for row in rows])
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/history/<int:history_id>", methods=["GET"])
def get_history_item(history_id):
    try:
        with sqlite3.connect(DATABASE) as conn:
            conn.row_factory = sqlite3.Row
            cursor = conn.cursor()
            cursor.execute('SELECT * FROM history WHERE id = ?', (history_id,))
            row = cursor.fetchone()
            if row:
                return jsonify(dict(row))
            return jsonify({"error": "Item not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    print(f"Eco AI Server starting on http://localhost:5000")
    app.run(debug=True, port=5000)
